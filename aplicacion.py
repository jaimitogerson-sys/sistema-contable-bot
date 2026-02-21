import os
import requests
import json
import time
import threading
from datetime import datetime
from flask import Flask, request
from google.oauth2 import service_account
from googleapiclient.discovery import build
import pandas as pd
import fitz  # PyMuPDF
import docx2txt
import pytesseract
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import openai

app = Flask(__name__)

# ================== VARIABLES ==================
TOKEN = os.getenv("TOKEN")                   # Token Telegram
CHAT_ID = os.getenv("CHAT_ID")               # Tu chat_id
FOLDER_ID = os.getenv("FOLDER_ID")           # Carpeta de Drive a monitorear
GOOGLE_CREDENTIALS = json.loads(os.getenv("GOOGLE_CREDENTIALS"))
OPENAI_KEY = os.getenv("OPENAI_KEY")         # API Key OpenAI

DESCARGA_FOLDER = "descargas"
RESULTADOS_FOLDER = "resultados"
os.makedirs(DESCARGA_FOLDER, exist_ok=True)
os.makedirs(RESULTADOS_FOLDER, exist_ok=True)

openai.api_key = OPENAI_KEY
archivos_vistos = set()

# ================== TELEGRAM ==================
def enviar_mensaje(chat_id, texto):
    url = f"https://api.telegram.org/bot{TOKEN}/sendMessage"
    data = {"chat_id": chat_id, "text": texto, "parse_mode": "Markdown"}
    requests.post(url, data=data)

# ================== GOOGLE DRIVE ==================
creds = service_account.Credentials.from_service_account_info(
    GOOGLE_CREDENTIALS, scopes=["https://www.googleapis.com/auth/drive"]
)
drive_service = build("drive", "v3", credentials=creds)

def descargar_archivo(file_id, nombre_local=None):
    if not nombre_local:
        nombre_local = os.path.join(DESCARGA_FOLDER, f"{file_id}")
    request_drive = drive_service.files().get_media(fileId=file_id)
    from googleapiclient.http import MediaIoBaseDownload
    import io
    fh = io.FileIO(nombre_local, 'wb')
    downloader = MediaIoBaseDownload(fh, request_drive)
    done = False
    while not done:
        status, done = downloader.next_chunk()
    return nombre_local

# ================== LECTURA DE ARCHIVOS ==================
def leer_excel(file_path):
    return pd.read_excel(file_path)

def leer_pdf(file_path):
    texto = ""
    doc = fitz.open(file_path)
    for pagina in doc:
        texto += pagina.get_text()
    return texto

def leer_word(file_path):
    return docx2txt.process(file_path)

def leer_imagen(file_path):
    return pytesseract.image_to_string(Image.open(file_path))

def extraer_datos(file_path):
    ext = file_path.split(".")[-1].lower()
    try:
        if ext in ["xls","xlsx"]:
            return leer_excel(file_path)
        elif ext == "pdf":
            return leer_pdf(file_path)
        elif ext in ["doc","docx"]:
            return leer_word(file_path)
        elif ext in ["png","jpg","jpeg"]:
            return leer_imagen(file_path)
        else:
            return f"No se puede procesar: {file_path}"
    except Exception as e:
        return f"Error leyendo {file_path}: {e}"

# ================== DRIVE FUNCIONES ==================
def buscar_carpeta(nombre, parent_id):
    query = f"mimeType='application/vnd.google-apps.folder' and name='{nombre}' and '{parent_id}' in parents and trashed=false"
    resultado = drive_service.files().list(q=query, fields="files(id, name)").execute()
    carpetas = resultado.get("files", [])
    return carpetas[0]["id"] if carpetas else None

def crear_carpeta_si_no_existe(nombre, parent_id):
    carpeta = buscar_carpeta(nombre, parent_id)
    if carpeta:
        return carpeta
    metadata = {"name": nombre, "mimeType": "application/vnd.google-apps.folder", "parents": [parent_id]}
    carpeta = drive_service.files().create(body=metadata, fields="id").execute()
    return carpeta["id"]

def crear_estructura_carpetas(ruta, folder_base):
    partes = ruta.split("/")
    parent_actual = folder_base
    for nombre in partes:
        parent_actual = crear_carpeta_si_no_existe(nombre, parent_actual)
    return parent_actual

def mover_archivo(file_id, carpeta_destino):
    archivo = drive_service.files().get(fileId=file_id, fields="parents").execute()
    padres_anteriores = ",".join(archivo.get("parents"))
    drive_service.files().update(fileId=file_id, addParents=carpeta_destino, removeParents=padres_anteriores, fields="id, parents").execute()

# ================== GENERAR RESULTADOS ==================
def generar_excel(dataframes, nombre="resultado.xlsx"):
    path = os.path.join(RESULTADOS_FOLDER, nombre)
    with pd.ExcelWriter(path) as writer:
        for sheet_name, df in dataframes.items():
            df.to_excel(writer, sheet_name=sheet_name, index=False)
    return path

def generar_presentacion(resumen, nombre="presentacion.pptx"):
    path = os.path.join(RESULTADOS_FOLDER, nombre)
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    left, top, width, height = Inches(0.5), Inches(1.5), Inches(9), Inches(5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    for r in resumen:
        p = tf.add_paragraph()
        p.text = str(r)
        p.font.size = Pt(14)
        p.alignment = PP_ALIGN.LEFT
    prs.save(path)
    return path

# ================== IA PARA INSTRUCCIONES ==================
def procesar_instruccion_ia(instruccion, archivos):
    archivos_info = {nombre:str(type(extraer_datos(ruta))) for nombre,ruta in archivos.items()}
    prompt = f"""
Tienes los archivos: {archivos_info}.
Instrucción: {instruccion}

Realiza la tarea indicada, analiza los archivos y genera resultados concretos. Indica si crear Excel, PPT u otros archivos.
Devuelve un resumen claro y un plan de acción.
"""
    response = openai.ChatCompletion.create(
        model="gpt-4",
        messages=[{"role":"user","content":prompt}],
        temperature=0
    )
    return response.choices[0].message.content

# ================== MONITOREO DRIVE ==================
def revisar_drive():
    global archivos_vistos
    while True:
        try:
            resultados = drive_service.files().list(q=f"'{FOLDER_ID}' in parents and trashed=false", fields="files(id, name)").execute()
            archivos = resultados.get("files", [])
            nuevos = [a for a in archivos if a["id"] not in archivos_vistos]

            if nuevos:
                for a in nuevos:
                    archivos_vistos.add(a["id"])
                nombres = "\n".join([f"📄 {a['name']}" for a in nuevos])
                enviar_mensaje(CHAT_ID, f"📥 *NUEVO ARCHIVO DETECTADO*\n\n{nombres}\n\n📂 Envía tu instrucción para procesarlos.")
            time.sleep(20)
        except Exception as e:
            print("Error revisando Drive:", e)
            time.sleep(30)

# ================== TELEGRAM WEBHOOK ==================
@app.route(f"/{TOKEN}", methods=["POST"])
def telegram_webhook():
    data = request.json
    if "message" in data:
        texto = data["message"].get("text", "").strip()
        if not texto:
            return "ok"
        inicio = time.time()
        try:
            # Crear carpetas
            if "crear carpeta" in texto.lower():
                ruta = texto.lower().replace("crear carpeta","").strip()
                crear_estructura_carpetas(ruta,FOLDER_ID)
                tiempo = round(time.time()-inicio,2)
                enviar_mensaje(CHAT_ID,f"✅ Carpeta creada: {ruta}\n⏱ Tiempo: {tiempo}s")
                return "ok"

            # Descargar archivos
            archivos_locales = {}
            for a_id in archivos_vistos:
                file_meta = drive_service.files().get(fileId=a_id, fields="name").execute()
                file_name = file_meta["name"]
                file_local = descargar_archivo(a_id, os.path.join(DESCARGA_FOLDER,file_name))
                archivos_locales[file_name] = file_local

            # Procesar instrucción con IA
            resultado_ia = procesar_instruccion_ia(texto, archivos_locales)
            tiempo_total = round(time.time()-inicio,2)
            enviar_mensaje(CHAT_ID, f"✅ Instrucción ejecutada:\n{texto}\n\n📄 Resultado:\n{resultado_ia}\n⏱ Tiempo total: {tiempo_total}s")
        except Exception as e:
            enviar_mensaje(CHAT_ID,f"❌ Error: {e}")
    return "ok"

# ================== HOME ==================
@app.route("/")
def home():
    return "Sistema IA Contable Activo"

# ================== HILO MONITOREO ==================
hilo = threading.Thread(target=revisar_drive)
hilo.daemon = True
hilo.start()

# ================== INICIO SERVIDOR ==================
if __name__ == "__main__":
    port = int(os.environ.get("PORT",10000))
    app.run(host="0.0.0.0",port=port)
